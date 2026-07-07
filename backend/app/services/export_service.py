"""Export concept artifacts (notes/slides/quiz) as downloadable Office files.

Notes and quiz render to .docx (python-docx, already a dependency); slides
render to .pptx with speaker notes (python-pptx). Builders read the same JSON
shapes the generators emit (generation_prompts.py schemas) and skip any
section that's absent, so partial content still exports.
"""
import io
import re

from fastapi import HTTPException
from supabase import Client

from app.services import generation_service as gen

MEDIA = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _safe_name(s: str) -> str:
    return re.sub(r"[^\w\- ]+", "", s or "").strip().replace(" ", "_") or "concept"


def _concept_name(db: Client, job_id: str, concept_id: str) -> str:
    res = (
        db.table("artifacts").select("content")
        .eq("job_id", job_id).eq("type", "topic_plan").limit(1).execute()
    )
    if res.data:
        for c in (res.data[0].get("content") or {}).get("concept_inventory", []):
            if c.get("concept_id") == concept_id:
                return c.get("concept_name") or concept_id
    return concept_id


def export_concept(db: Client, job_id: str, topic_id: str, concept_id: str, artifact_type: str):
    """Return (bytes, filename, media_type) for the artifact, or 404 if not generated."""
    row = gen._concept_row(db, topic_id, concept_id, artifact_type)
    content = (row or {}).get("content")
    if not content:
        raise HTTPException(status_code=404, detail="Not generated yet")
    name = _concept_name(db, job_id, concept_id)
    if artifact_type == "slides":
        data, ext = _slides_pptx(content, name), "pptx"
    elif artifact_type == "quiz":
        data, ext = _quiz_docx(content, name), "docx"
    else:
        data, ext = _notes_docx(content, name), "docx"
    suffix = {"student_notes": "notes", "slides": "slides", "quiz": "quiz"}.get(artifact_type, artifact_type)
    return data, f"{_safe_name(name)}_{suffix}.{ext}", MEDIA[ext]


# ── builders ──────────────────────────────────────────────────────────────────

def _notes_docx(content: dict, name: str) -> bytes:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(name, level=0)
    core = (content or {}).get("core") or {}
    opening = (content or {}).get("opening") or {}
    closing = (content or {}).get("closing") or {}
    osec = opening.get("sections") or {}
    csec = closing.get("sections") or {}

    def h(t): doc.add_heading(t, level=1)

    def p(t):
        # Notes fields may be prose (legacy), arrays of points, or
        # {core, elaboration} objects (structured schema).
        if not t:
            return
        if isinstance(t, dict):
            p(t.get("core"))
            bullets(t.get("elaboration"))
        elif isinstance(t, list):
            bullets(t)
        else:
            doc.add_paragraph(str(t))

    def bullets(items):
        for it in items or []:
            doc.add_paragraph(str(it), style="List Bullet")

    ps = osec.get("problem_statement") or {}
    intro = (osec.get("introduction") or {}).get("narrative_intro")
    if ps.get("scenario") or intro:
        h("Why this matters")
        p(ps.get("scenario")); p(ps.get("gap_statement")); p(intro)

    outs = (osec.get("topic_overview") or {}).get("outcomes_checklist") or []
    if outs:
        h("Learning outcomes")
        bullets([o.get("statement") for o in outs if o.get("statement")])

    cc = core.get("core_concept") or {}
    if cc.get("formal_definition"):
        h("Definition")
        p(cc["formal_definition"]); p(cc.get("mental_model_analogy"))

    dd = core.get("deep_dive") or {}
    mech = dd.get("architecture_and_mechanism") or {}
    if mech.get("explanation"):
        h("Architecture & mechanism"); p(mech["explanation"])

    code = dd.get("code_or_formalization") or {}
    if code.get("applicable") and code.get("content"):
        h(f"Code ({code.get('language_or_system') or code.get('type') or ''})".replace(" ()", ""))
        cpara = doc.add_paragraph(str(code["content"]))
        for run in cpara.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(9)
        p(code.get("explanation"))

    pu = core.get("practical_understanding") or {}
    if pu.get("worked_example"):
        h("Worked example"); p(pu["worked_example"])
    if pu.get("advantages"):
        h("Advantages"); bullets(pu["advantages"])
    if pu.get("disadvantages"):
        h("Trade-offs"); bullets(pu["disadvantages"])
    if pu.get("applications"):
        h("Real-world applications"); bullets(pu["applications"])

    mistakes = csec.get("common_mistakes") or []
    if mistakes:
        h("Common mistakes")
        for m in mistakes:
            wrong = m.get("wrong_way") or m.get("mistake") or ""
            why = m.get("why_it_fails") or m.get("correct_approach") or ""
            doc.add_paragraph(f"{wrong} — {why}".strip(" —"), style="List Bullet")

    rev = csec.get("revision_section") or {}
    if rev.get("key_takeaways"):
        h("Key takeaways"); bullets(rev["key_takeaways"])
    if rev.get("important_definitions"):
        h("Important definitions")
        for d in rev["important_definitions"]:
            doc.add_paragraph(f"{d.get('term')}: {d.get('definition')}", style="List Bullet")

    glossary = (csec.get("glossary_section") or {}).get("terms") or []
    if glossary:
        h("Glossary")
        for t in glossary:
            doc.add_paragraph(
                f"{t.get('term')}: {t.get('formal_definition') or t.get('simple_explanation') or ''}",
                style="List Bullet")

    cards = (csec.get("flashcard_section") or {}).get("cards") or []
    if cards:
        h("Flashcards")
        for c in cards:
            doc.add_paragraph(f"Q: {c.get('front')} — A: {c.get('back')}", style="List Bullet")

    pq = csec.get("practice_questions") or {}
    groups = [("Easy", pq.get("easy") or []), ("Medium", pq.get("medium") or []),
              ("Hard", pq.get("hard") or [])]
    if any(qs for _, qs in groups):
        h("Practice questions")
        for label, qs in groups:
            for i, q in enumerate(qs, 1):
                doc.add_paragraph(f"[{label} {i}] {q.get('question', '')}")
        doc.add_heading("Answers", level=2)
        for label, qs in groups:
            for i, q in enumerate(qs, 1):
                if q.get("answer_explanation"):
                    doc.add_paragraph(f"[{label} {i}] {q['answer_explanation']}")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _quiz_docx(content: dict, name: str) -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading(f"{name} — Quiz", level=0)
    letters = "ABCDEF"
    mcq = (content or {}).get("mcq") or []
    sa = (content or {}).get("short_answer") or []

    if mcq:
        doc.add_heading("Multiple choice", level=1)
        for i, q in enumerate(mcq, 1):
            doc.add_paragraph(f"Q{i}. {q.get('question', '')}")
            for oi, o in enumerate(q.get("options") or []):
                label = letters[oi] if oi < len(letters) else str(oi + 1)
                doc.add_paragraph(f"    {label}. {o}")
    if sa:
        doc.add_heading("Short answer", level=1)
        for i, q in enumerate(sa, 1):
            doc.add_paragraph(f"{i}. {q.get('question', '')}")

    if mcq or sa:
        doc.add_page_break()
        doc.add_heading("Answer key", level=1)
        for i, q in enumerate(mcq, 1):
            ai = q.get("answer_index")
            letter = letters[ai] if isinstance(ai, int) and 0 <= ai < len(letters) else "?"
            line = f"Q{i}: {letter}"
            if q.get("explanation"):
                line += f" — {q['explanation']}"
            doc.add_paragraph(line)
        for i, q in enumerate(sa, 1):
            if q.get("model_answer"):
                doc.add_paragraph(f"Short answer {i}: {q['model_answer']}")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _slides_pptx(content: dict, name: str) -> bytes:
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=501, detail="python-pptx is not installed on the server")

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = name

    body_layout = prs.slide_layouts[1]  # Title and Content
    for s in (content or {}).get("slides") or []:
        sl = prs.slides.add_slide(body_layout)
        sl.shapes.title.text = str(s.get("title") or "")
        lines: list[str] = []
        # Layout-specific fields, flattened in reading order.
        if s.get("definition_core"):
            lines.append(str(s["definition_core"]))
        for sec in s.get("sections") or []:
            if isinstance(sec, dict):
                lines.append(f"{sec.get('heading', '')}:")
                lines.extend(f"  {b}" for b in (sec.get("bullets") or []))
        for t in s.get("terms") or []:
            if isinstance(t, dict):
                lines.append(f"{t.get('term', '')} — {t.get('definition', '')}")
        if s.get("left_bullets") or s.get("right_bullets"):
            lines.append(f"{s.get('left_heading') or 'Advantages'}:")
            lines.extend(f"  {b}" for b in (s.get("left_bullets") or []))
            lines.append(f"{s.get('right_heading') or 'Limitations'}:")
            lines.extend(f"  {b}" for b in (s.get("right_bullets") or []))
        lines.extend(str(b) for b in (s.get("body_blocks") or []))
        if s.get("myth"):
            lines.insert(0, f"Myth: {s['myth']}")
        if s.get("reality"):
            lines.insert(1 if s.get("myth") else 0, f"Reality: {s['reality']}")
        code = (s.get("code") or {}).get("content")
        if code:
            lines.extend(str(code).splitlines())
        if s.get("takeaway"):
            lines.append(f"★ {s['takeaway']}")
        if lines:
            tf = sl.placeholders[1].text_frame
            tf.text = lines[0]
            for b in lines[1:]:
                tf.add_paragraph().text = b
        if s.get("speaker_notes"):
            sl.notes_slide.notes_text_frame.text = str(s["speaker_notes"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
